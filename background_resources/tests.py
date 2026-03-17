import os
import shutil
import time
import zipfile
from pathlib import Path
from datetime import datetime

from django.test import TestCase
from django.test.utils import override_settings
from django.conf import settings
from django.core.files.uploadedfile import SimpleUploadedFile

from background_resources.models import (
    Document, 
    ReadingStrategy, 
    PromptStrategy, 
    RegexStrategy, 
    AbbreviationsReadingStrategy,
    RAGChunk,
    StrategyChunkUsage,
    RAGQueryLog
)
from llm_api.apps import service_registry

# Libraries for generating test files
from reportlab.pdfgen import canvas
from docx import Document as DocxDocument
from pptx import Presentation

from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings


TEST_BASE_DIR = Path(settings.BASE_DIR) / "test_data"
TEST_VECTOR_STORE = TEST_BASE_DIR / "vector_store"
TEST_CHUNK_STORE = TEST_BASE_DIR / "chunk_store"
TEST_FILES_DIR = TEST_BASE_DIR / "files"
TEST_RESULTS_DIR = Path(settings.BASE_DIR) / "test_results"


class TestFileGenerator:
    """Helper class to generate valid files for testing ingestion."""

    @staticmethod
    def create_txt(path, content="This is a simple text file for testing."):
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return path

    @staticmethod
    def create_pdf(path, content="This is a PDF document for testing."):
        c = canvas.Canvas(str(path))
        c.drawString(100, 750, content)
        c.drawString(100, 730, "Page 1 Content")
        c.showPage()
        c.drawString(100, 750, "Page 2 Content")
        c.save()
        return path

    @staticmethod
    def create_docx(path, content="This is a DOCX document for testing."):
        doc = DocxDocument()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph(content)
        doc.add_page_break()
        doc.add_paragraph("Second page content.")
        doc.save(path)
        return path

    @staticmethod
    def create_pptx(path, content="This is a PPTX presentation for testing."):
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = content
        prs.save(path)
        return path

    @staticmethod
    def create_zip_corpus(path, content="<html><body><p>This is a ZIP corpus content block.</p></body></html>"):
        with zipfile.ZipFile(path, 'w') as zf:
            zf.writestr("index.html", content)
            # Add a nested file to ensure recursive loading works
            zf.writestr("subdir/nested.html", "<html><body><p>Nested ZIP content.</p></body></html>")
        return path

class BackgroundResourcesIntegrationTest(TestCase):
    """
    Full Integration Test Suite.
    NO MOCKS. Uses real AI Service, NLP Service, and RAG Service.
    """

    @classmethod
    def setUpClass(cls):
        # Manually apply settings override for the duration of the class
        # This ensures setUpClass (and the service registry) sees the test paths
        cls.settings_override = override_settings(
            VECTOR_STORE=TEST_VECTOR_STORE,
            CHUNK_STORE=TEST_CHUNK_STORE,
            FILES=TEST_FILES_DIR,
            MEDIA_ROOT=TEST_FILES_DIR
        )
        cls.settings_override.enable()
        super().setUpClass()
        # Create test directories
        os.makedirs(TEST_VECTOR_STORE, exist_ok=True)
        os.makedirs(TEST_CHUNK_STORE, exist_ok=True)
        os.makedirs(TEST_FILES_DIR, exist_ok=True)
        os.makedirs(TEST_RESULTS_DIR, exist_ok=True)

        print("\n>>> 🚀 INITIALIZING REAL SERVICES (This may take time) <<<")

        cls.ai_service = service_registry['ai_service']
        cls.nlp_service = service_registry['nlp_service']
        cls.rag_service = service_registry['rag_service']

        # Ensure Spacy is loaded for Abbreviation tests
        cls.nlp_service.get_abbreviation_model()

    @classmethod
    def tearDownClass(cls):
        # Cleanup test data
        if os.path.exists(TEST_BASE_DIR):
            shutil.rmtree(TEST_BASE_DIR)
        super().tearDownClass()
        cls.settings_override.disable()

    def setUp(self):
        # Reset the RAG service DB for a clean state per test
        embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        self.rag_service.db = FAISS.from_texts(["init"], embeddings)
        init_id = list(self.rag_service.db.docstore._dict.keys())[0]
        self.rag_service.db.delete([init_id])
        self.start_time = 0

    def _create_document_obj(self, file_path):
        """Helper to create a Django Document object from a file path."""
        with open(file_path, 'rb') as f:
            django_file = SimpleUploadedFile(
                name=os.path.basename(file_path),
                content=f.read(),
                content_type='application/octet-stream'
            )
        
        doc = Document.objects.create(
            title=f"Test {os.path.basename(file_path)}",
            file=django_file,
            chunk_size=500, # Default for tests
            chunk_overlap=50
        )
        return doc

    def _generate_report(self, test_name, duration):
        """Generates a markdown report of what is in the vector store."""
        report_path = TEST_RESULTS_DIR / f"{test_name}_report.md"
        
        vs_docs = self.rag_service.db.docstore._dict
        
        with open(report_path, "w", encoding="utf-8") as f:
            f.write(f"# Test Report: {test_name}\n\n")
            f.write(f"**Date:** {datetime.now()}\n")
            f.write(f"**Duration:** {duration:.4f} seconds\n")
            
            f.write("## Vector Store Contents\n")
            f.write(f"Total Chunks: {len(vs_docs)}\n\n")
            
            for chunk_id, doc in vs_docs.items():
                # Try to fetch full content from store if available
                store_id = doc.metadata.get("chunk_id") or doc.metadata.get("doc_id")
                full_content = "[Content in Vector Only]"
                if store_id:
                    try:
                        stored_doc = self.rag_service.store.mget([store_id])[0]
                        if stored_doc:
                            full_content = stored_doc.page_content
                    except:
                        pass

                f.write(f"### Chunk ID: {chunk_id}\n")
                f.write(f"- **Metadata:** {doc.metadata}\n")
                f.write(f"- **Vector Content:** {doc.page_content[:200]}...\n")
                f.write(f"- **Store Content:** {full_content[:200]}...\n")
                f.write("---\n")
        
        print(f"\nReport generated: {report_path}")

    def test_1_default_ingestion_and_overrides(self):
        """
        GOAL: Ingest document with default strategy, then ingest same document
        with overridden chunking strategy. Confirm chunks are new and different sizes.
        """
        print("\n>>> Test 1: Default Ingestion & Overrides")
        self.start_time = time.perf_counter()

        # 1. Create a long text file to ensure multiple chunks
        long_text = "This is a sentence. " * 500 # ~10k characters
        txt_path = TestFileGenerator.create_txt(TEST_FILES_DIR / "long_doc.txt", long_text)
        
        doc = self._create_document_obj(txt_path)

        # 2. Default Reading Strategy
        print("   -> Running Default Reading Strategy...")
        default_strategy = ReadingStrategy.objects.create(
            document=doc,
            strategy_description="Default Reading"
        )
        default_strategy.read_document(self.rag_service)

        # Verify Default Chunks
        default_usages = default_strategy.usages.all()
        count_default = default_usages.count()
        print(f"      Default Strategy generated {count_default} chunks.")
        self.assertTrue(count_default > 0)

        # 3. Override Reading Strategy (Small Chunks)
        print("   -> Running Override Reading Strategy (Small Chunks)...")
        small_strategy = ReadingStrategy.objects.create(
            document=doc,
            strategy_description="Small Chunks",
            chunk_size_override=100, # Force very small chunks
            chunk_overlap_override=10
        )
        small_strategy.read_document(self.rag_service)

        # Verify Override Chunks
        small_usages = small_strategy.usages.all()
        count_small = small_usages.count()
        print(f"      Small Strategy generated {count_small} chunks.")
        
        # Assertion: Smaller chunk size should result in MORE chunks for the same text
        self.assertTrue(count_small > count_default, 
                        f"Override should produce more chunks ({count_small}) than default ({count_default})")

        self._generate_report("test_1_overrides", time.perf_counter() - self.start_time)

    def test_2_subclass_strategies(self):
        """
        GOAL: Create and exercise PromptStrategy, RegexStrategy, and AbbreviationsReadingStrategy.
        """
        print("\n>>> Test 2: Subclass Strategies")
        self.start_time = time.perf_counter()

        # --- A. Regex Strategy (Glossary) ---
        print("   -> Testing RegexStrategy...")
        glossary_content = """
        GLOSSARY
        Alpha-Code: The first code in the sequence.
        Beta-Code - The second code in the sequence.
        """
        glossary_path = TestFileGenerator.create_txt(TEST_FILES_DIR / "glossary.txt", glossary_content)
        glossary_doc = self._create_document_obj(glossary_path)
        
        regex_strat = RegexStrategy.objects.create(
            document=glossary_doc,
            strategy_description="Glossary Extraction",
            regex=r"(.+?)(?:\s*:(?!\/\/)\s*|\s*-\s+)(.+)" # Capture "Alpha-Code" and definition
        )
        regex_strat.apply_strategy(self.rag_service)
        
        # Verify: Check if "Alpha-Code" is in the vector store (as a term)
        # and the definition is in the store.
        vs_docs = self.rag_service.db.docstore._dict
        found_term = any("Alpha-Code" in d.page_content for d in vs_docs.values())
        self.assertTrue(found_term, "RegexStrategy should index the Term")

        # --- B. Abbreviations Strategy ---
        print("   -> Testing AbbreviationsReadingStrategy...")
        abrv_content = "The National Aeronautics and Space Administration (NASA) explores space."
        abrv_path = TestFileGenerator.create_txt(TEST_FILES_DIR / "abrv.txt", abrv_content)
        abrv_doc = self._create_document_obj(abrv_path)

        abrv_strat = AbbreviationsReadingStrategy.objects.create(
            document=abrv_doc,
            strategy_description="Abbreviation Extraction"
        )
        abrv_strat.apply_strategy(self.rag_service)

        # Verify: Spacy should find NASA
        vs_docs = self.rag_service.db.docstore._dict
        found_nasa = any("NASA" in d.page_content for d in vs_docs.values())
        self.assertTrue(found_nasa, "AbbreviationsStrategy should index 'NASA'")

        # --- C. Prompt Strategy (Real AI) ---
        print("   -> Testing PromptStrategy (Calling AI Service)...")
        # Keep text short to save time/tokens
        prompt_content = "The quick brown fox jumps over the lazy dog."
        prompt_path = TestFileGenerator.create_txt(TEST_FILES_DIR / "prompt.txt", prompt_content)
        prompt_doc = self._create_document_obj(prompt_path)

        prompt_strat = PromptStrategy.objects.create(
            document=prompt_doc,
            strategy_description="AI Summary",
            prompt="Summarize this text, describing the contents in a long expression, then further refining it to a concise description in a short expression, and then choosing the most appropriate keywords:"
        )
        prompt_strat.apply_strategy(self.rag_service)

        # Verify: We expect *some* output in the vector store that isn't the original text
        # The PromptStrategy indexes the *summary*, not the raw text.
        # We check if we have a chunk associated with this strategy.
        prompt_usages = prompt_strat.usages.all()
        self.assertTrue(prompt_usages.exists(), "PromptStrategy should generate chunks")

        self._generate_report("test_2_strategies", time.perf_counter() - self.start_time)

    def test_3_retrieval_and_stats(self):
        """
        GOAL: Confirm that Chunk hit_count and last_accessed fields are correctly updated.
        """
        print("\n>>> Test 3: Retrieval & Stats")
        self.start_time = time.perf_counter()

        # 1. Ingest a document with a unique keyword
        unique_keyword = "Xylophone"
        content = f"The {unique_keyword} is a musical instrument."
        path = TestFileGenerator.create_txt(TEST_FILES_DIR / "stats.txt", content)
        doc = self._create_document_obj(path)
        
        strategy = ReadingStrategy.objects.create(document=doc, strategy_description="Stats Test")
        strategy.read_document(self.rag_service)

        # 2. Verify initial state (hit_count = 0)
        usage = strategy.usages.first()
        self.assertEqual(usage.chunk.hit_count, 0)
        initial_access = usage.chunk.last_accessed

        # 3. Perform Retrieval
        print(f"   -> Querying for '{unique_keyword}'...")
        results = self.rag_service.get_context(unique_keyword, k=1)
        
        self.assertTrue(len(results) > 0, "Should retrieve the document")
        self.assertIn(unique_keyword, results[0].page_content)

        # 4. Verify Stats Update
        usage.refresh_from_db()
        usage.chunk.refresh_from_db()
        print(f"   -> Chunk Hit Count: {usage.chunk.hit_count}")
        self.assertEqual(usage.chunk.hit_count, 1, "Hit count should increment after retrieval")
        self.assertNotEqual(usage.chunk.last_accessed, initial_access, "Last accessed timestamp should update")

        # 5. Verify Query Log
        log = RAGQueryLog.objects.last()
        self.assertEqual(log.query_text, unique_keyword)
        print(f"   -> Query Logged: {log.query_text}")

        self._generate_report("test_3_stats", time.perf_counter() - self.start_time)

    def test_4_file_loaders(self):
        """
        GOAL: Test ingestion of PDF, DOCX, PPTX, and ZIP (Corpora) files using the default ReadingStrategy.
        """
        print("\n>>> Test 4: File Loaders (PDF, DOCX, PPTX, ZIP)")
        self.start_time = time.perf_counter()

        # 1. Generate Files
        pdf_path = TestFileGenerator.create_pdf(TEST_FILES_DIR / "test.pdf", "This is a PDF content block.")
        docx_path = TestFileGenerator.create_docx(TEST_FILES_DIR / "test.docx", "This is a DOCX content block.")
        pptx_path = TestFileGenerator.create_pptx(TEST_FILES_DIR / "test.pptx", "This is a PPTX content block.")
        zip_path = TestFileGenerator.create_zip_corpus(TEST_FILES_DIR / "test.zip", "<html><body><p>This is a ZIP corpus content block.</p></body></html>")

        files = [pdf_path, docx_path, pptx_path, zip_path]
        
        for path in files:
            print(f"   -> Processing {os.path.basename(path)}...")
            doc = self._create_document_obj(path)
            strategy = ReadingStrategy.objects.create(
                document=doc,
                strategy_description=f"Default {os.path.basename(path)}"
            )
            strategy.read_document(self.rag_service)
            
            # Verify chunks
            usages = strategy.usages.all()
            self.assertTrue(usages.exists(), f"Should generate chunks for {os.path.basename(path)}")
            
            # Verify content in vector store
            chunk_id = usages.first().chunk.chunk_id
            stored_doc = self.rag_service.store.mget([chunk_id])[0]
            self.assertIsNotNone(stored_doc)
            
            if "pdf" in str(path):
                self.assertIn("PDF content", stored_doc.page_content)
            elif "docx" in str(path):
                self.assertIn("DOCX content", stored_doc.page_content)
            elif "pptx" in str(path):
                self.assertIn("PPTX content", stored_doc.page_content)
            elif "zip" in str(path):
                self.assertIn("ZIP corpus content", stored_doc.page_content)

        # Verify Corpus Cleanup (Specific to ZIP)
        zip_doc = Document.objects.get(title__contains="test.zip")
        corpus_path = Path(settings.MEDIA_ROOT) / "corpora" / zip_doc.indexed_hash
        self.assertTrue(corpus_path.exists(), "Corpus folder should exist after ingestion")
        
        print(f"   -> Verifying cleanup for {corpus_path}...")
        zip_doc.delete()
        self.assertFalse(corpus_path.exists(), "Corpus folder should be deleted after document deletion")

        self._generate_report("test_4_loaders", time.perf_counter() - self.start_time)

    def test_5_regex_precision(self):
        """
        GOAL: Verify that RegexStrategy extracts ONLY the definition line,
        even if the source file has tricky line endings (like \r) that might
        cause the regex to match the whole chunk greedily.
        """
        print("\n>>> Test 5: Regex Precision (Line Endings)")
        self.start_time = time.perf_counter()

        # 1. Create a file with \r line endings (Classic Mac / some CSV exports)
        # If the splitter fails to split this, it becomes one long line.
        content = "Term1: Definition One.\rTerm2: Definition Two.\rTerm3: Definition Three."
        path = TestFileGenerator.create_txt(TEST_FILES_DIR / "tricky_glossary.txt", content)
        doc = self._create_document_obj(path)

        # 2. Run Regex Strategy
        strategy = RegexStrategy.objects.create(
            document=doc,
            strategy_description="Tricky Glossary",
            regex=r"(.+?)(?:\s*:(?!\/\/)\s*|\s*-\s+)(.+)"
        )
        strategy.apply_strategy(self.rag_service)

        # 3. Verify Store Content
        # We expect "Definition One." NOT "Definition One.\rTerm2: Definition Two..."
        usages = strategy.usages.all()
        self.assertTrue(usages.exists())
        
        # Check the first extracted definition
        # We need to find the chunk corresponding to "Term1"
        # Since we can't easily query the ByteStore by content, we check the Chunk text_content cache
        # (which we added in a previous step for Admin visibility)
        
        # We iterate to find the specific definition
        found_clean_def = False
        for usage in usages:
            if "Definition One" in usage.chunk.text_content:
                print(f"   -> Checking content: {repr(usage.chunk.text_content)}")
                if "Definition Two" not in usage.chunk.text_content:
                    found_clean_def = True
                    break
        
        self.assertTrue(found_clean_def, "Regex should have stopped at the line break, avoiding greedy capture of subsequent terms.")

        self._generate_report("test_5_regex", time.perf_counter() - self.start_time)

    def test_6_deletion_robustness(self):
        """
        GOAL: Verify that chunks are removed from the index/store ONLY when
        no higher-order object relies on them.
        """
        print("\n>>> Test 6: Deletion Robustness")
        self.start_time = time.perf_counter()

        # 1. Create Document & Default Strategy
        content = "Shared content for robustness test."
        path = TestFileGenerator.create_txt(TEST_FILES_DIR / "robustness.txt", content)
        doc = self._create_document_obj(path)

        # Ensure ingestion happened for default strategy (created by signal)
        default_strat = ReadingStrategy.objects.get(document=doc, strategy_description="Default Chunking")
        default_strat.read_document(self.rag_service)
        
        default_usages = default_strat.usages.all()
        self.assertTrue(default_usages.exists())
        chunk_id = default_usages.first().chunk.chunk_id
        
        # Verify in store
        self.assertIsNotNone(self.rag_service.store.mget([chunk_id])[0], "Chunk should be in store")

        # 2. Create a second strategy reusing these chunks
        prompt_strat = PromptStrategy.objects.create(
            document=doc,
            strategy_description="Dependent Strategy",
            prompt="Summarize"
        )
        prompt_strat.apply_strategy(self.rag_service)
        
        # Check DB count - Should be 2 (Default + Prompt reference)
        self.assertEqual(StrategyChunkUsage.objects.filter(chunk__chunk_id=chunk_id).count(), 2, "Should have 2 references to this chunk ID")

        # 3. Delete the PromptStrategy
        print("   -> Deleting PromptStrategy...")
        prompt_strat.delete()
        
        # Verify Default chunks still exist in DB
        self.assertEqual(StrategyChunkUsage.objects.filter(chunk__chunk_id=chunk_id).count(), 1)

        # Verify Store Content EXISTENCE (The Critical Check)
        in_store = self.rag_service.store.mget([chunk_id])[0]
        self.assertIsNotNone(in_store, "Chunk should STILL be in store after deleting one reference")

        # 4. Delete the Document (cascades to Default Strategy)
        print("   -> Deleting Document...")
        doc.delete()
        
        # Verify everything gone
        self.assertEqual(StrategyChunkUsage.objects.filter(chunk__chunk_id=chunk_id).count(), 0)
        in_store_final = self.rag_service.store.mget([chunk_id])[0]
        self.assertIsNone(in_store_final, "Chunk should be gone from store")

        self._generate_report("test_6_robustness", time.perf_counter() - self.start_time)
